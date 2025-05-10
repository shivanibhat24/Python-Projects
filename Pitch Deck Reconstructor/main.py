"""
Pitch Deck Reconstructor

This script transforms a meeting transcript into a professional startup pitch deck.
It uses Whisper for audio-to-text (if needed), LLaMA 3 for content analysis,
LangChain for orchestration, and python-pptx for presentation generation.

Dependencies:
- openai
- langchain
- python-pptx
- huggingface_hub
- torch
- transformers

Date: May 4, 2025
"""

import os
import re
import json
import argparse
from typing import List, Dict, Any, Optional, Tuple
import torch
from pathlib import Path

# LLM-related imports
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.llms import HuggingFacePipeline
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline

# Whisper for audio transcription (if needed)
from transformers import WhisperProcessor, WhisperForConditionalGeneration

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set up argparse
parser = argparse.ArgumentParser(description='Generate a pitch deck from a transcript')
parser.add_argument('--input', type=str, required=True, help='Path to transcript text file or audio file')
parser.add_argument('--output', type=str, default='pitch_deck.pptx', help='Output presentation filename')
parser.add_argument('--template', type=str, help='Optional PowerPoint template to use')
parser.add_argument('--company', type=str, help='Company name (optional)')

# Define color schemes
COLOR_SCHEMES = {
    "modern": {
        "primary": RGBColor(0, 112, 192),    # Blue
        "secondary": RGBColor(255, 192, 0),  # Gold
        "text": RGBColor(68, 68, 68),        # Dark Gray
        "accent": RGBColor(112, 173, 71),    # Green
    },
    "tech": {
        "primary": RGBColor(77, 77, 255),    # Bright Blue
        "secondary": RGBColor(61, 61, 61),   # Dark Gray
        "text": RGBColor(33, 33, 33),        # Nearly Black
        "accent": RGBColor(0, 204, 153),     # Teal
    },
    "corporate": {
        "primary": RGBColor(31, 73, 125),    # Navy Blue
        "secondary": RGBColor(192, 0, 0),    # Dark Red
        "text": RGBColor(51, 51, 51),        # Dark Gray
        "accent": RGBColor(128, 128, 128),   # Medium Gray
    }
}

class LlamaModel:
    """Wrapper for LLaMA 3 model to handle content generation"""
    
    def __init__(self, model_id="meta-llama/Meta-Llama-3-8B-Instruct"):
        """Initialize the LLaMA 3 model"""
        print(f"Loading LLaMA 3 model: {model_id}")
        self.tokenizer = AutoTokenizer.from_pretrained(model_id)
        self.model = AutoModelForCausalLM.from_pretrained(
            model_id,
            torch_dtype=torch.float16,
            device_map="auto",
        )
        self.pipe = pipeline(
            "text-generation",
            model=self.model,
            tokenizer=self.tokenizer,
            max_new_tokens=2048,
            temperature=0.7,
            top_p=0.95,
            repetition_penalty=1.15
        )
        self.llm = HuggingFacePipeline(pipeline=self.pipe)
    
    def get_langchain_llm(self):
        """Return the LangChain-compatible LLM"""
        return self.llm

class TranscriptProcessor:
    """Process audio or text transcripts"""
    
    def __init__(self):
        """Initialize the transcript processor"""
        self.whisper_model = None
        self.whisper_processor = None
    
    def load_whisper(self):
        """Load Whisper model for audio transcription"""
        print("Loading Whisper model for audio transcription...")
        self.whisper_processor = WhisperProcessor.from_pretrained("openai/whisper-large-v3")
        self.whisper_model = WhisperForConditionalGeneration.from_pretrained("openai/whisper-large-v3")
        if torch.cuda.is_available():
            self.whisper_model = self.whisper_model.to("cuda")
        print("Whisper model loaded successfully")
    
    def process_input(self, input_path: str) -> str:
        """Process the input file (either text or audio) and return the transcript"""
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        file_ext = os.path.splitext(input_path)[1].lower()
        
        # Handle text files
        if file_ext in ['.txt', '.md', '.doc', '.docx']:
            print(f"Processing text file: {input_path}")
            with open(input_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # Handle audio files
        elif file_ext in ['.mp3', '.wav', '.m4a', '.ogg']:
            print(f"Processing audio file: {input_path}")
            if self.whisper_model is None:
                self.load_whisper()
            
            import librosa
            audio, sr = librosa.load(input_path, sr=16000)
            
            # Process with Whisper
            input_features = self.whisper_processor(audio, sampling_rate=sr, return_tensors="pt").input_features
            if torch.cuda.is_available():
                input_features = input_features.to("cuda")
            
            predicted_ids = self.whisper_model.generate(input_features)
            transcript = self.whisper_processor.batch_decode(predicted_ids, skip_special_tokens=True)[0]
            
            # Save transcript for reference
            transcript_path = os.path.splitext(input_path)[0] + "_transcript.txt"
            with open(transcript_path, 'w', encoding='utf-8') as f:
                f.write(transcript)
            
            print(f"Transcript saved to {transcript_path}")
            return transcript
        
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

class PitchDeckAnalyzer:
    """Analyze transcript and extract pitch deck content"""
    
    def __init__(self, llm):
        """Initialize with a language model"""
        self.llm = llm
    
    def extract_company_info(self, transcript: str, company_name: Optional[str] = None) -> Dict[str, Any]:
        """Extract basic company information from the transcript"""
        prompt = PromptTemplate(
            input_variables=["transcript", "company_hint"],
            template="""
            Analyze this startup pitch transcript and extract the company name, tagline/slogan, 
            and industry/sector. If you can't find specific information, make a reasonable guess 
            based on context.
            
            Company name hint (use if provided): {company_hint}
            
            Transcript:
            {transcript}
            
            Return your analysis as JSON with these keys:
            - company_name: The name of the company
            - tagline: A brief slogan or tagline
            - industry: The primary industry or sector
            - color_scheme: Suggested color scheme (choose one: modern, tech, corporate)
            """
        )
        
        chain = LLMChain(llm=self.llm, prompt=prompt)
        response = chain.run(transcript=transcript, company_hint=company_name or "")
        
        # Clean up response to extract JSON
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            response = json_match.group(0)
        
        try:
            company_info = json.loads(response)
            # Ensure we have all required keys
            for key in ["company_name", "tagline", "industry", "color_scheme"]:
                if key not in company_info:
                    company_info[key] = ""
            return company_info
        except json.JSONDecodeError:
            # Fallback if JSON parsing fails
            print("Warning: Failed to parse JSON response from LLM")
            return {
                "company_name": company_name or "Startup",
                "tagline": "Innovative Solutions",
                "industry": "Technology",
                "color_scheme": "modern"
            }
    
    def analyze_transcript(self, transcript: str, company_info: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze the full transcript and extract content for each slide"""
        prompt = PromptTemplate(
            input_variables=["transcript", "company_name", "tagline", "industry"],
            template="""
            You are an expert pitch deck consultant. Analyze this startup pitch transcript for {company_name} 
            in the {industry} industry with the tagline "{tagline}".
            
            Extract content for a comprehensive pitch deck with these slides:
            1. Title Slide (company name, tagline, presenter)
            2. Problem Statement (what problem is being solved)
            3. Solution (how the company solves this problem)
            4. Market Opportunity (market size, trends, growth potential)
            5. Product/Service (details about what they offer)
            6. Business Model (how they make money)
            7. Competitive Landscape (competitors and advantages)
            8. Traction & Metrics (current progress, customers, revenue)
            9. Team (key team members and background)
            10. Financial Projections (revenue forecasts, funding needs)
            11. Go-to-Market Strategy (how they'll acquire customers)
            12. Call to Action (what they're asking for)
            
            Transcript:
            {transcript}
            
            For each slide, extract:
            1. A concise title
            2. 2-4 bullet points of key information (complete sentences)
            3. Any specific numbers, metrics or quotes to highlight
            
            Return your analysis as JSON with a key for each slide type and structured content.
            """
        )
        
        chain = LLMChain(llm=self.llm, prompt=prompt)
        response = chain.run(
            transcript=transcript,
            company_name=company_info["company_name"],
            tagline=company_info["tagline"],
            industry=company_info["industry"]
        )
        
        # Clean up response to extract JSON
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            response = json_match.group(0)
        
        try:
            return json.loads(response)
        except json.JSONDecodeError:
            print("Warning: Failed to parse JSON response from LLM")
            # Return a minimal structure if parsing fails
            return {
                "title_slide": {
                    "title": company_info["company_name"],
                    "subtitle": company_info["tagline"],
                    "bullets": []
                },
                "problem": {
                    "title": "Problem Statement",
                    "bullets": ["Problem identified in the market"]
                },
                "solution": {
                    "title": "Our Solution",
                    "bullets": ["How we solve the problem"]
                }
            }

class PowerPointGenerator:
    """Generate PowerPoint presentation from analyzed content"""
    
    def __init__(self, template_path: Optional[str] = None):
        """Initialize with optional template"""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            print(f"Using template from {template_path}")
        else:
            self.prs = Presentation()
            print("Using default PowerPoint template")
        
        # Create slide layouts reference
        self.layouts = {
            'title': 0,  # Title slide layout
            'title_content': 1,  # Title and content layout
            'section': 2,  # Section header layout
            'two_content': 3,  # Two content layout
            'comparison': 4,  # Comparison layout
            'blank': 5,  # Blank layout
        }
    
    def apply_theme(self, color_scheme_name: str):
        """Apply a color scheme to the presentation"""
        if color_scheme_name not in COLOR_SCHEMES:
            color_scheme_name = "modern"
        
        self.colors = COLOR_SCHEMES[color_scheme_name]
        print(f"Applied {color_scheme_name} color scheme")
    
    def add_title_slide(self, company_name: str, tagline: str):
        """Add the title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layouts['title']])
        
        # Add title and subtitle
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = company_name
        subtitle.text = tagline
        
        # Style the title
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(44)
                run.font.color.rgb = self.colors["primary"]
        
        # Style the subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.italic = True
                run.font.size = Pt(24)
                run.font.color.rgb = self.colors["secondary"]
        
        return slide
    
    def add_content_slide(self, title: str, bullets: List[str], highlight: Optional[str] = None):
        """Add a content slide with title and bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layouts['title_content']])
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Style the title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(32)
                run.font.color.rgb = self.colors["primary"]
        
        # Add bullets
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = self.colors["text"]
        
        # Add highlight if provided
        if highlight:
            left = Inches(0.5)
            top = Inches(5)
            width = Inches(9)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = highlight
            p.alignment = PP_ALIGN.CENTER
            
            for run in p.runs:
                run.font.italic = True
                run.font.size = Pt(20)
                run.font.color.rgb = self.colors["accent"]
        
        return slide
    
    def generate_presentation(self, content: Dict[str, Any], company_info: Dict[str, Any]) -> str:
        """Generate the full presentation from analyzed content"""
        # Apply color scheme
        color_scheme = company_info.get("color_scheme", "modern")
        self.apply_theme(color_scheme)
        
        # Title slide
        self.add_title_slide(
            company_info["company_name"],
            company_info["tagline"]
        )
        
        # Standard slides based on content keys
        slide_keys = [
            "problem", "solution", "market_opportunity", "product", "business_model",
            "competition", "traction", "team", "financials", "go_to_market", "call_to_action"
        ]
        
        for key in slide_keys:
            if key in content:
                slide_content = content[key]
                # Extract bullets as a list
                bullets = slide_content.get("bullets", [])
                if isinstance(bullets, str):
                    # Split string into list if necessary
                    bullets = [b.strip() for b in bullets.split('\n') if b.strip()]
                
                highlight = slide_content.get("highlight", "")
                
                self.add_content_slide(
                    slide_content.get("title", key.replace("_", " ").title()),
                    bullets,
                    highlight
                )
        
        return self.save_presentation(company_info["company_name"])
    
    def save_presentation(self, company_name: str) -> str:
        """Save the presentation to disk"""
        # Clean company name for filename
        safe_name = re.sub(r'[^\w\-_]', '_', company_name)
        filename = f"{safe_name}_pitch_deck.pptx"
        
        self.prs.save(filename)
        print(f"Presentation saved as {filename}")
        return filename

def main():
    """Main execution function"""
    args = parser.parse_args()
    
    # Process input file (transcript or audio)
    processor = TranscriptProcessor()
    try:
        transcript = processor.process_input(args.input)
        print(f"Successfully processed input: {len(transcript)} characters")
    except Exception as e:
        print(f"Error processing input: {e}")
        return
    
    # Initialize LLaMA 3 model
    try:
        llama = LlamaModel()
        llm = llama.get_langchain_llm()
        print("Successfully initialized LLaMA 3 model")
    except Exception as e:
        print(f"Error initializing LLaMA 3: {e}")
        return
    
    # Analyze transcript
    analyzer = PitchDeckAnalyzer(llm)
    
    # Extract company info
    company_info = analyzer.extract_company_info(transcript, args.company)
    print(f"Company identified: {company_info['company_name']} - {company_info['industry']}")
    
    # Analyze full transcript
    deck_content = analyzer.analyze_transcript(transcript, company_info)
    print(f"Successfully analyzed transcript and extracted content for {len(deck_content)} slides")
    
    # Generate presentation
    ppt_gen = PowerPointGenerator(args.template)
    output_file = ppt_gen.generate_presentation(deck_content, company_info)
    
    print(f"✅ Pitch deck generation complete!")
    print(f"📊 Output file: {output_file}")

if __name__ == "__main__":
    main()
